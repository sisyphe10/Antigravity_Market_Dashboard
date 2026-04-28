# -*- coding: utf-8 -*-
"""NH 목표전환형 2호 세미나 PPT 빌드 (v2 - fixed)"""
import copy, sys, io, os, json
from pptx import Presentation
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

BASE = r'C:/Users/user/Antigravity_Market_Dashboard'
SRC = os.path.join(BASE, 'PPT', '260313_NH_목표전환형_랩_다이내믹밸류.pptx')
DST = SRC  # overwrite

# ====================== helpers ======================
def delete_slide_by_idx(prs, idx0):
    xml_slides = prs.slides._sldIdLst
    slide_elements = list(xml_slides)
    xml_slides.remove(slide_elements[idx0])

def duplicate_slide(prs, source_idx0):
    """Append deep copy of source slide to presentation end."""
    source = prs.slides[source_idx0]
    slide_layout = source.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)
    # copy rels first so deepcopied shapes retain working r:embed/r:link
    for rel_id, rel in list(source.part.rels.items()):
        if rel.reltype.endswith('/notesSlide'):
            continue
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            try:
                new_slide.part.relate_to(rel._target, rel.reltype)
            except Exception:
                pass
    for shp in source.shapes:
        el = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.append(el)
    return new_slide, len(prs.slides) - 1

def text_replace_many(slide, mapping):
    """Para-level replacement. Collapses split runs into first run."""
    shapes = list(slide.shapes)
    for shp in shapes:
        if shp.has_text_frame:
            _replace_in_tf(shp.text_frame, mapping)
        elif shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    _replace_in_tf(cell.text_frame, mapping)

def _replace_in_tf(tf, mapping):
    """Single-pass, non-cascading replacement using regex alternation."""
    import re
    keys = [k for k in mapping.keys() if k]
    keys.sort(key=len, reverse=True)  # longest first so "AB" beats "A"
    if not keys:
        return
    pattern = re.compile('|'.join(re.escape(k) for k in keys))
    for para in tf.paragraphs:
        full = ''.join(run.text for run in para.runs)
        if not full.strip():
            continue
        new_full = pattern.sub(lambda m: mapping.get(m.group(0), m.group(0)), full)
        if new_full != full:
            if para.runs:
                first = para.runs[0]
                first.text = new_full
                for run in para.runs[1:]:
                    run._r.getparent().remove(run._r)
            else:
                run = para.add_run()
                run.text = new_full

def set_cell_text(cell, text):
    """Replace a table cell's text, preserving first run style."""
    tf = cell.text_frame
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = str(text)
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            run = p.add_run()
            run.text = str(text)

def get_first_table(slide):
    for shp in slide.shapes:
        if shp.has_table:
            return shp.table
    return None

def get_all_tables(slide):
    return [shp.table for shp in slide.shapes if shp.has_table]


# ====================== main build ======================
def main():
    print(f'Load: {SRC}')
    prs = Presentation(SRC)
    n = len(prs.slides)
    print(f'Slides: {n}')

    # Do NOT delete slides upfront — that causes partname reuse collisions.
    # Instead, keep all 46 originals, duplicate templates (which get unique slide47+ partnames),
    # then build final _sldIdLst containing only the slides we want in the order we want.
    # Originals we didn't reference become orphans in the package (harmless).

    def idx_of(orig1):
        return orig1 - 1  # 0-based index in original 46-slide deck

    # Duplicate templates (appended to end, partnames slide47.xml ~ slide55.xml)
    dup = {}
    for key, src in [
        ('headline',  9),   # big chart layout
        ('ai',       14),   # text-heavy layout
        ('multipolar',14),
        ('kculture', 14),
        ('diagnosis',14),
        ('pf_frame', 30),   # 4-grid layout
        ('qa',       14),
        ('three_values', 39),  # 5-box layout
        ('track',    42),   # performance table layout
    ]:
        _, new_idx = duplicate_slide(prs, idx_of(src))
        dup[key] = new_idx
        print(f'  dup {src} → idx {new_idx} [{key}]')

    # Current 0-based index for each logical slide
    cur = {
        'cover':         idx_of(1),
        'co_divider':    idx_of(2),
        'co_overview':   idx_of(3),
        'co_mgmt':       idx_of(4),
        'co_org':        idx_of(5),
        'mk_divider':    idx_of(7),
        'mk_headline':   dup['headline'],
        'mk_3axes':      idx_of(8),
        'mk_ai':         dup['ai'],
        'mk_multipolar': dup['multipolar'],
        'mk_law':        idx_of(14),
        'mk_kculture':   dup['kculture'],
        'mk_diagnosis':  dup['diagnosis'],
        'mk_pf_frame':   dup['pf_frame'],
        'mk_end_check':  idx_of(30),
        'mk_qa':         dup['qa'],
        'ph_divider':    idx_of(31),
        'ph_value_esg':  idx_of(32),
        'ph_engagement': idx_of(33),
        'ph_process':    idx_of(34),
        'ph_arbitrage':  idx_of(36),
        'pr_sisyphus':   idx_of(38),
        'pr_overview':   idx_of(40),
        'pr_strategy':   idx_of(41),
        'pr_three_values': dup['three_values'],
        'pr_track':      dup['track'],
        'pr_live':       idx_of(42),
        'pr_portfolio':  idx_of(43),
        'contact':       idx_of(46),
    }

    final = [
        'cover','co_divider','co_overview','co_mgmt','co_org',
        'mk_divider','mk_headline','mk_3axes','mk_ai','mk_multipolar',
        'mk_law','mk_kculture','mk_diagnosis','mk_pf_frame','mk_end_check','mk_qa',
        'ph_divider','ph_value_esg','ph_engagement','ph_process','ph_arbitrage',
        'pr_sisyphus','pr_overview','pr_strategy','pr_three_values','pr_track',
        'pr_live','pr_portfolio','contact',
    ]
    print(f'Final: {len(final)} slides')

    # Reorder
    xml_slides = prs.slides._sldIdLst
    slide_els = list(xml_slides)
    for el in slide_els:
        xml_slides.remove(el)
    for k in final:
        xml_slides.append(slide_els[cur[k]])

    # Updated dict with final indices
    S = {k: prs.slides[i] for i, k in enumerate(final)}

    apply_updates(S)

    print(f'Save: {DST}')
    prs.save(DST)
    print('Done.')


def apply_updates(S):
    # ========== cover ==========
    text_replace_many(S['cover'], {
        'Dynamic Value': 'NH 목표전환형 2호',
        'Mar. 2026': 'Apr. 2026',
    })

    # ========== mk_headline (dup from slide 9) ==========
    # Orig 9 texts: '외국인의 시각에서 본 KOSPI',
    # '시장 전망 : 2026년의 상승 동력은 Multiple Rerating',
    # "KOSPI/USD의 '신고가'는 2026년 들어서야 시작",
    # '(출처: 한국거래소, 라이프자산운용)',
    # disclaimer footer
    text_replace_many(S['mk_headline'], {
        '외국인의 시각에서 본 KOSPI': '코스피 6,500 — Premium KOSPI의 현실화',
        '시장 전망 : 2026년의 상승 동력은 Multiple Rerating': '시황 및 전망',
        "KOSPI/USD의 '신고가'는 2026년 들어서야 시작": 'Discount의 정상화, Premium KOSPI의 원년',
        "KOSPI/USD의 ‘신고가’는 2026년 들어서야 시작": 'Discount의 정상화, Premium KOSPI의 원년',
        '(출처: 한국거래소, 라이프자산운용)': '· 2026-04-22 KOSPI 6,417.93 (YTD +52.3% · 1Y +158.1%)\n· 3차 상법 개정 완료 · 배당분리과세 시행 임박 → 지배구조 정상화\n· P/B 리레이팅 진행, 12M P/E는 여전히 박스권 하단 → 추가 여지 존재',
    })

    # ========== mk_3axes (orig 8) ==========
    text_replace_many(S['mk_3axes'], {
        '2026년 상반기 코스피 전망': '시장을 이끄는 3가지 축 + α',
        '2026년 상반기 코스피 시장: Premium Korea': '',
        'Discount의 정상화, Premium KOSPI의 원년': '',
        '정부의 시장친화 정책, 제도 변화': 'AI · 데이터센터',
        '자본시장 환경 개선': '반도체 · 변압기 · 에너지 · 800V DC',
        '개인, 외국인, 기관 수급의3박자 유동성 랠리': '산업 · 안보 · 무역 다극화',
        '2025년 가격 상승에도 불구하고여전히 매력적인 멀티플': '조선 · 방산 · 원전 · 신재생',
        '자본 시장으로의 머니 무브를 원하는 정부,친시장 정책 드라이브 본격화': '상법 개정',
        '상법 개정, 배당소득 분리 과세 등 제도 개선으로 긍정적 여건 변화': '국내 시장 멀티플 상승',
        '한국 시장 리스크 감소 (할인율 감소)': '(+α) K-culture · K-pop · K-cosmetic · 인바운드 · 내수',
        '아직까지 신중한 리테일 분위기': '글로벌 AI 인프라 투자 사이클의 핵심 수혜',
        '외국인 25년 누적 순매수: -4.6조원': '인접 산업 전반 동반 수혜 확산',
        '연기금의 국내 주식시장 투자 비율 상향, 퇴직연금 기금화 등': '공급망 재편의 구조적 수혜국',
        'P = PER * EPS작년의 주가지수 상승은 멀티플 상향보다 EPS 증가에 기반': '중장기 한국 제조·에너지 섹터 재평가',
        '여전히 낮은 PER, PBR이 글로벌 투자자 유입 요인으로 작용': '주주환원 강화·지배구조 개선',
    })

    # ========== mk_ai (dup of 14) ==========
    # Orig 14 canonical texts - keep very specific to avoid cross-replacement
    apply_law_template_replacements(S['mk_ai'], {
        'title': '축① AI → 데이터센터 — 반도체 · 에너지 · 인프라 동반 수혜',
        'header': '시황 및 전망 | 축① AI · 데이터센터',
        'left_header': '글로벌 AI 인프라 투자 확장 국면',
        'left_sub': '반도체·전력·인프라 섹터 동반 수혜',
        'left_b1': '· 글로벌 LLM 토큰 사용량 폭증: 텍스트 → 이미지 → 비디오 → AI Agent',
        'left_b2': '· HBM 수요 급증, DRAM Wafer start 정체 → 수급 구조적 타이트 지속 예상',
        'right_header': '인접 수혜 산업',
        'right_items': [
            '· 변압기 · 전력기기 — AI 데이터센터 전력 수요 급증',
            '· 에너지(LNG · 원전 · 태양광) — 전력 인프라 확충',
            '· 800V DC — 데이터센터 고효율 전력 공급 표준화',
            '· 해상 데이터센터 등 신 테마 확산',
        ],
        'progress_header': '핵심 관전 포인트',
        'progress_items': [
            '· 반도체 Peer 대비 극단적 저평가 → 리레이팅 진행',
            '· 이익 추정치 상향 지속 → 멀티플 정상화 동력',
            '· AI Capex 사이클 → 한국 반도체·장비·소재 구조적 수혜',
        ],
    })

    # ========== mk_multipolar ==========
    apply_law_template_replacements(S['mk_multipolar'], {
        'title': '축② 다극화 시대 — 공급망 재편의 핵심 수혜국',
        'header': '시황 및 전망 | 축② 다극화',
        'left_header': '산업 · 안보 · 무역 다극화의 구조적 수혜',
        'left_sub': '공급망 재편 → 한국 전략적 포지션 확대',
        'left_b1': '· 미-중 갈등 장기화 · 공급망 분절 → 중간 포지션 한국 경쟁력 부각',
        'left_b2': '· 중장기 수혜 섹터: 조선 · 방산 · 원전 · 신재생',
        'right_header': '최근 지정학 흐름',
        'right_items': [
            '· 중동 휴전 국면 진입 → 유가 안정',
            '· 유가 안정 → 재건 · 건설 · 인프라 섹터로 수혜 확산',
            '· 원전 · 태양광 등 에너지 전환 섹터 부각',
            '· AI 인프라 전력 수요와 결합되며 구조적 성장 동력 강화',
        ],
        'progress_header': '투자 시사점',
        'progress_items': [
            '· 다극화는 단기 이슈가 아닌 중장기 구조적 변화',
            '· 한국은 에너지·안보·산업 밸류체인의 핵심 거점',
            '· 조선 · 방산 · 원전 · 신재생 섹터의 중장기 재평가 기대',
        ],
    })

    # ========== mk_law (orig 14 in place) — 상법 + 배당분리과세 merged ==========
    apply_law_template_replacements(S['mk_law'], {
        'title': '축③ 제도 변화 — 상법 개정 + 배당분리과세',
        'header': '시황 및 전망 | 축③ 제도 변화',
        'left_header': '상법 개정안 국회 통과',
        'left_sub': '기업 지배구조 투명성 및 주주 권리 보호 강화',
        'left_b1': '· 기업의 투명한 의사결정 · 소액주주 권익 보호 · 주주간 이해상충 완화',
        'left_b2': '· 코리아 디스카운트 요인 해소로 외국인 투자 유입, 증시 재평가 기대',
        'right_header': '상법 개정안 주요 내용',
        'right_items': [
            '· 이사의 충실의무 확대 (회사 → 회사 및 주주)',
            '· 상장회사 전자주주총회 의무화',
            '· 사외이사 독립성 강화',
            '· 감사위원 분리 선출',
            '· 집중투표제 의무화',
            '· 자사주 소각 의무화 (3차 개정 통과)',
        ],
        'progress_header': '진행 현황 및 배당분리과세',
        'progress_items': [
            '· 1차 시행 2025.7.22 / 2차 통과 2025.8.25 / 3차 통과 2026.2.25',
            '· 배당분리과세 법안 통과 (2025.12.2), 2026년 배당분부터 적용',
            '· 배당성향 40%↑ 또는 25%↑ & 증가율 10%↑ 상장사 → 배당금 종합소득 제외',
        ],
    })

    # ========== mk_kculture ==========
    apply_law_template_replacements(S['mk_kculture'], {
        'title': '+α K-culture · 내수 낙수효과의 시작',
        'header': '시황 및 전망 | +α K-culture · 내수',
        'left_header': 'K-culture 모멘텀 확산',
        'left_sub': '글로벌 수요 기반 강세 지속',
        'left_b1': '· K-pop · K-cosmetic — 글로벌 수요 확대 지속',
        'left_b2': '· 인바운드 관광 회복 → 내수 및 서비스 산업 활력',
        'right_header': '내수 낙수효과 확인 구간',
        'right_items': [
            '· 반도체 · IT 대형사 성과급 본격화 → 내수 회복 시그널',
            '· 일부 백화점 매출 YoY +20% 이상 성장 — 내수 턴 포인트 초입',
            '· 강세장 후반 내수 · 소비재 · 유통 수혜 확산 기대',
            '',
        ],
        'progress_header': '투자 시사점',
        'progress_items': [
            '· 수출·제조 중심 랠리에서 내수로 온기 확산',
            '· K-브랜드 경쟁력 기반 구조적 성장 동력 유지',
            '· 소비재 · 유통 섹터의 실적·주가 회복 기대',
        ],
    })

    # ========== mk_diagnosis ==========
    apply_law_template_replacements(S['mk_diagnosis'], {
        'title': '현대 한국의 최전성기 — 운용의 묘가 필요한 구간',
        'header': '시황 및 전망 | 현 국면 진단',
        'left_header': '중장기 구조적 강세 여건',
        'left_sub': '자생력 높은 산업 구조와 해외 진출',
        'left_b1': '· 자생 가능한 산업 구조 — 반도체·조선·자동차·에너지·문화 전반',
        'left_b2': '· 해외 진출 본격화 국면 진입 → 글로벌 시장 점유율 확대 전망',
        'right_header': '단기 변동성 특성',
        'right_items': [
            '· 지수 전고점 갱신 구간마다 패시브 차익실현에 따른 단기 변동 반복 가능',
            '· 뉴스 · 재료에 대한 시장 반응의 속도와 진폭 확대',
            '· 단순 매수·보유 대비 탄력적 운용의 상대적 효율 상승',
            '',
        ],
        'progress_header': '운용 시사점',
        'progress_items': [
            '· 널리 알려진 재료라도 탄력적 매매 역량이 성과의 핵심',
            '· 변동성 · MDD 적극 관리가 장기 수익률에 결정적 영향',
            '· 목표전환형 랩의 효용이 극대화되는 시장 국면',
        ],
    })

    # ========== mk_pf_frame (dup of 30, 4-grid) ==========
    apply_risk_template_replacements(S['mk_pf_frame'], {
        'header': '시황 및 전망 | 포트폴리오 프레임',
        'title': '어떤 섹터에 주목할 것인가',
        'subtitle': '테마 리더 · 대형 섹터 순환매 · 코스닥 대형주 3가지 프레임 제시',
        'items': [
            '① 테마 리더 + 신 테마',
            'AI · 다극화 해당 산업 및 결합 신테마 (해상 데이터센터 · 800V DC · 원전 · 태양광)',
            '② 대형 섹터 순환매',
            '패시브 시대 매매 효율적인 대형주 · 자동차 · 조선 등 기간조정 섹터 순환매',
            '③ 코스닥 대형주',
            '실적·LO 레퍼런스 확실한 대형주 · 바이오 정리 · 2차전지 턴어라운드',
        ],
    })

    # ========== mk_end_check (orig 30 in place) ==========
    apply_risk_template_replacements(S['mk_end_check'], {
        'header': '시황 및 전망 | 모니터링 포인트',
        'title': '강세장이 언제 끝날까',
        'subtitle': '강세장 지속 여부를 판단하는 4가지 모니터링 포인트',
        'items': [
            '① AI 사이클 — LLM 성능 정체 / 독점 / 과금 한계',
            '② 밸류에이션 — 실적 둔화 속 멀티플 팽창',
            '③ 지정학 — 미-중 관계 봉합 · 다극화 정착 · 정면충돌',
            '④ K-모멘텀 — 글로벌 투자자의 한국 관심도 저하',
            '네 가지 지표를 선제적 모니터링',
            '국면별 포지션 조절 예정',
        ],
    })

    # ========== mk_qa (dup of 14) ==========
    apply_law_template_replacements(S['mk_qa'], {
        'title': '자주 묻는 질문',
        'header': '시황 및 전망 | Q & A',
        'left_header': 'Q1. 증시는 좋은데 체감 경기는 아직 어렵지 않은가?',
        'left_sub': '낙수효과가 확인되는 구간',
        'left_b1': '· 풀뿌리 경기의 장기간 부진은 사실이었음',
        'left_b2': '· 반도체·IT 대형사 성과급 본격화로 낙수효과 확인 · 일부 백화점 매출 YoY +20% 이상',
        'right_header': 'Q2. 코스닥은 언제 코스피를 따라잡을 것인가?',
        'right_items': [
            '· 연초 전망 대비 가장 지연된 구간이나 기저 축적 진행 중',
            '· 바이오 대형주 우려 이슈는 정리 국면',
            '· 2차전지: 중국 내 구조조정 · 유가 안정 · ESS · BBU 수요로 턴어라운드 기대',
            '',
        ],
        'progress_header': '선별 접근 원칙',
        'progress_items': [
            '· 실적 · LO(Long Only) 레퍼런스 확실한 대형주 중심',
            '· 단기 뉴스보다 구조적 모멘텀 기반 종목 선별',
            '· 변동성 관리 기반 포트폴리오 운용으로 진입 타이밍 리스크 완화',
        ],
    })

    # ========== pr_overview (orig 40) ==========
    text_replace_many(S['pr_overview'], {
        'NH-라이프 목표전환형 랩': 'NH 목표전환형 2호',
        '3.20(금) - 3.24(화)': '별도 고지',
        '3.25(수)': '별도 고지',
        'WRAP 상품 개요': 'NH 목표전환형 2호 상품 개요',
    })

    # ========== pr_strategy (orig 41) — minor touch ==========
    text_replace_many(S['pr_strategy'], {
        'NH-라이프 목표전환형 랩 운용 전략': 'NH 목표전환형 2호 운용 전략',
    })

    # ========== pr_three_values (dup of 39, 5-box → 3-box adapted) ==========
    # Blanking 04/05 circle numbers and 5th box label to de-clutter
    text_replace_many(S['pr_three_values'], {
        '운용 전략 개요': '목표전환형 랩의 3가지 가치',
        'NH-라이프 목표전환형 랩': 'NH 목표전환형 2호',
        '                  투자 철학': '              ① 변동성·MDD 관리',
        '                  운용 전략': '              ② 수익 확정·재운용',
        '                  리스크 관리': '              ③ 신규자금 마중물',
        '                  관심 섹터': '',
        '                  운용 목표': '                  목표전환형 랩이 제공하는 세 가지 투자 가치',
        '변동성/MDD 관리 하에 목표 수익률 조기 달성': '하우스 철학 기반 변동성·MDD 적극 관리',
        '가치 투자, 책임 투자, 기대감 차익거래': '하락 후 장기 미체결 리스크 완화',
        '반도체, 원전, 엔터, 로봇, 인바운드, 태양광, 바이오 등': '목표 수익률 조기 달성 시 자동 전환 → 수익 확정 후 재운용',
        '탄력적인 현금 운용, 섹터·종목 분산 투자': '신규·추가 자금의 마중물 → 진입 타이밍 고민 경감',
        '코스피-코스닥 바벨 전략(주도주-소외주)': '변동성 관리 × 수익 확정 × 타이밍 솔루션의 3박자',
        '04': '',
        '05': '',
    })

    # ========== pr_track (dup of 42) — 트랙 레코드 3건 ==========
    # Cell-level table update
    slide = S['pr_track']
    tables = get_all_tables(slide)
    if tables:
        t = tables[0]
        # Header row
        set_cell_text(t.cell(0, 0), '구분')
        set_cell_text(t.cell(0, 1), 'DB 1차')
        set_cell_text(t.cell(0, 2), 'DB 2차')
        set_cell_text(t.cell(0, 3), 'NH 1호')
        set_cell_text(t.cell(0, 4), '공통')
        # Row 1: 기준가
        set_cell_text(t.cell(1, 0), '최종 기준가')
        set_cell_text(t.cell(1, 1), '1,079.98')
        set_cell_text(t.cell(1, 2), '1,080.75')
        set_cell_text(t.cell(1, 3), '1,101.58')
        set_cell_text(t.cell(1, 4), '—')
        # Row 2: 수익률
        set_cell_text(t.cell(2, 0), '전환 수익률')
        set_cell_text(t.cell(2, 1), '+8.00%')
        set_cell_text(t.cell(2, 2), '+8.08%')
        set_cell_text(t.cell(2, 3), '+10.16%')
        set_cell_text(t.cell(2, 4), '목표 달성')
        # Row 3: 거래일
        set_cell_text(t.cell(3, 0), '소요 기간')
        set_cell_text(t.cell(3, 1), '10거래일')
        set_cell_text(t.cell(3, 2), '22거래일')
        set_cell_text(t.cell(3, 3), '16거래일')
        set_cell_text(t.cell(3, 4), '전원 달성')
    text_replace_many(slide, {
        '성과 추이: NH-라이프 개방형 랩 (2023년 3월 출시)': '목표전환형 랩 트랙 레코드 — 선행 3개 상품 모두 목표 수익률 달성',
        '시장 대비 낮은 MDD로 안정적인 수익 실현': '선행 3건 전원 목표 달성 · 3월 조정 구간 극복',
        'NH 개방형 랩은 2023년 3월 설정, 2026년 수익률은 3월 10일까지의 YTD 기준.': '기준가는 최종 전환 시점 기준. 과거 성과는 미래 수익을 보장하지 않습니다.',
        '시장이 하락으로 마감했던 2024년에도 절대 수익률(+) 달성': 'NH 목표전환형 2호는 동일한 운용 철학·프로세스로 설계·운용 예정',
    })

    # ========== pr_live (orig 42) — 운용중 개방형 랩 성과 ==========
    slide = S['pr_live']
    tables = get_all_tables(slide)
    if tables:
        t = tables[0]
        # Original columns: 수익률 (%) | 2023년 | 2024년 | 2025년 | 2026년
        # New: 수익률 (%) | 1M | 3M | YTD | 1Y
        set_cell_text(t.cell(0, 0), '수익률 (%)')
        set_cell_text(t.cell(0, 1), '1M')
        set_cell_text(t.cell(0, 2), '3M')
        set_cell_text(t.cell(0, 3), 'YTD')
        set_cell_text(t.cell(0, 4), '1Y')
        # Row 1: WRAP = 개방형 랩
        set_cell_text(t.cell(1, 0), '개방형 랩')
        set_cell_text(t.cell(1, 1), '8.0')
        set_cell_text(t.cell(1, 2), '40.9')
        set_cell_text(t.cell(1, 3), '62.5')
        set_cell_text(t.cell(1, 4), '144.5')
        # Row 2: KOSPI
        set_cell_text(t.cell(2, 0), 'KOSPI')
        set_cell_text(t.cell(2, 1), '11.0')
        set_cell_text(t.cell(2, 2), '29.6')
        set_cell_text(t.cell(2, 3), '52.3')
        set_cell_text(t.cell(2, 4), '158.1')
        # Row 3: KOSDAQ
        set_cell_text(t.cell(3, 0), 'KOSDAQ')
        set_cell_text(t.cell(3, 1), '1.7')
        set_cell_text(t.cell(3, 2), '21.7')
        set_cell_text(t.cell(3, 3), '27.6')
        set_cell_text(t.cell(3, 4), '64.9')
    text_replace_many(slide, {
        '성과 추이: NH-라이프 개방형 랩 (2023년 3월 출시)': '참고 — 운용 중 개방형 랩 성과 (2026-04-22 기준)',
        '시장 대비 낮은 MDD로 안정적인 수익 실현': '동일한 운용 철학·프로세스로 중장기 성과 축적',
        'NH 개방형 랩은 2023년 3월 설정, 2026년 수익률은 3월 10일까지의 YTD 기준.': '개방형 랩은 2024년 5월 설정. YTD는 연초 이후 기준. 기간별 수익률은 2026-04-22 기준.',
        '시장이 하락으로 마감했던 2024년에도 절대 수익률(+) 달성': '장기 누적 기준 KOSPI 대비 견조한 알파 시현',
    })

    # ========== pr_portfolio (orig 43) — 포트폴리오 상위 5종목 공개 ==========
    slide = S['pr_portfolio']
    pf_path = os.path.join(BASE, 'portfolio_data.json')
    with open(pf_path, 'r', encoding='utf-8') as f:
        pf = json.load(f)
    holdings = list(pf.values())[0] if pf else []
    top = sorted(holdings, key=lambda x: -x.get('weight', 0))[:5]
    tables = get_all_tables(slide)
    if tables:
        t = tables[0]
        for i, h in enumerate(top[:5], 1):
            set_cell_text(t.cell(i, 1), h.get('sector', ''))
            set_cell_text(t.cell(i, 2), str(h.get('code', '')))
            set_cell_text(t.cell(i, 3), h.get('name', ''))
            set_cell_text(t.cell(i, 4), str(int(h.get('weight', 0))))
    text_replace_many(slide, {
        '성과 추이: NH-라이프 개방형 랩 (2023년 3월 출시)': '참고 — 운용 중 개방형 랩 포트폴리오 (상위 5종목)',
        '포트폴리오 구성': '포트폴리오 구성 (2026-04-22 기준)',
    })


# ------ 14번 레이아웃(상법 스타일) 업데이트 공통 ------
def apply_law_template_replacements(slide, d):
    """
    Slide 14 layout positions (paragraph-level anchors):
    - '지배 구조 개선을 법 개정을 통해 명문화 했다는 것은 매우 강력한 시그널' → title (top)
    - '시장 전망: 1. Bottom Level-Up: 상법 개정 및 배당소득 분리과세' → header breadcrumb
    - '상법 개정안 국회 통과' → left header (big)
    - '기업 지배구조 투명성 및 주주 권리 보호 강화' → left sub
    - '기업의 투명한 의사결정, 소액주주 권익 보호, 주주간 이해상충 문제 완화 등' → left bullet 1
    - '코리아 디스카운트 요인 해소로 외국인 투자 유입, 증시 재평가 기대' → left bullet 2
    - '상법 개정안 주요 내용' → right header
    - '이사의 충실의무 확대 ...' → right item 1
    - '상장회사 전자주주총회 의무화' → right item 2
    - '사외이사 독립성 강화' → right item 3
    - '감사위원 분리 선출 및 독립성 강화' → right item 4
    - '집중투표제 의무화' → right item 5 (we omit)
    - '자사주 소각 의무화 법안(3차)' → right item 6 (we omit)
    - '진행 현황' → progress header
    - '1차 상법 개정 시행 (2025.7.22)' → progress item 1
    - '2차 상법 개정 본회의 통과 (2025.8.25)' → progress item 2
    - '3차 상법 개정 본회의 통과 (2026.2.25)' → progress item 3
    """
    ri = list(d.get('right_items', []))
    while len(ri) < 6:
        ri.append('')
    pi = list(d.get('progress_items', []))
    while len(pi) < 3:
        pi.append('')
    mapping = {
        '지배 구조 개선을 법 개정을 통해 명문화 했다는 것은 매우 강력한 시그널': d.get('title', ''),
        '시장 전망: 1. Bottom Level-Up: 상법 개정 및 배당소득 분리과세': d.get('header', ''),
        '상법 개정안 국회 통과': d.get('left_header', ''),
        '기업 지배구조 투명성 및 주주 권리 보호 강화': d.get('left_sub', ''),
        '기업의 투명한 의사결정, 소액주주 권익 보호, 주주간 이해상충 문제 완화 등': d.get('left_b1', ''),
        '코리아 디스카운트 요인 해소로 외국인 투자 유입, 증시 재평가 기대': d.get('left_b2', ''),
        '(출처: 한국경제)': d.get('source', '(출처: 라이프자산운용)'),
        '상법 개정안 주요 내용': d.get('right_header', ''),
        '이사의 충실의무 확대 (기존 회사에서 [회사 및 주주]로 확대)': ri[0],
        '상장회사 전자주주총회 의무화': ri[1],
        '사외이사 독립성 강화': ri[2],
        '감사위원 분리 선출 및 독립성 강화': ri[3],
        '집중투표제 의무화': ri[4],
        '자사주 소각 의무화 법안(3차)': ri[5],
        '진행 현황': d.get('progress_header', ''),
        '1차 상법 개정 시행 (2025.7.22)': pi[0],
        '2차 상법 개정 본회의 통과 (2025.8.25)': pi[1],
        '3차 상법 개정 본회의 통과 (2026.2.25)': pi[2],
    }
    text_replace_many(slide, mapping)


def apply_risk_template_replacements(slide, d):
    """Slide 30 (리스크) layout — 4-grid + subtitle."""
    items = d.get('items', [''] * 6)
    mapping = {
        "모두가 예상하는 ' 2026년 상고하저', 그러나 역발상이 아닌 정발상이 필요한 시점": d.get('subtitle', ''),
        "모두가 예상하는 ‘ 2026년 상고하저’, 그러나 역발상이 아닌 정발상이 필요한 시점": d.get('subtitle', ''),
        '시장 전망: 리스크 요인': d.get('header', ''),
        '리스크 요인': d.get('title', ''),
        'AI 버블(?) 붕괴': items[0] if len(items) > 0 else '',
        '시장참여자들의 차익 실현': items[1] if len(items) > 1 else '',
        '한국 6월 지방선거': items[2] if len(items) > 2 else '',
        '미국 11월 중간선거': items[3] if len(items) > 3 else '',
        '원화 약세, 외국인 자금 유출': items[4] if len(items) > 4 else '',
        '지정학적 리스크 (미국-이란 전쟁 등)': items[5] if len(items) > 5 else '',
    }
    text_replace_many(slide, mapping)


if __name__ == '__main__':
    main()
